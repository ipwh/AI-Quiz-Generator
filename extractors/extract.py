# extractors/extract.py
# ---------------------------------------------------------
# 目標：香港中學 AI 題目生成器的統一抽取器（文字/圖片/Vision/OCR）
#
# 修正重點（保留所有功能）：
# P0：統一唯一 data URL helper（bytes_to_data_url）
# P1：extract_images_for_llm_ocr 改呼叫共用 helper（不重複實作）
# P2：OCR 品質閾值調低至 0.15；可讀字符集加入理科常用符號；亂碼比例判斷
# P3：繁體中文 OCR 優化（預處理、參數、後處理）
# P4（新增，Streamlit Cloud / 校本實務向）：
#    - 掃描 PDF OCR 加入頁數上限與分辨率參數，避免 OCR 全書拖死
#    - PDF 渲染優先用 dpi（若支援），fallback 才用 Matrix
#    - OCR 語言包缺失自動降級（chi_tra/chi_sim 不存在就移除）
#    - PSM fallback（結果太差會自動再試）
#    - 二值化 threshold 改為 Otsu（更適合不同掃描底色），可 fallback 固定值
# ---------------------------------------------------------

import io
import re
import base64
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation

# OCR (optional)
try:
    import pytesseract
    from PIL import Image, ImageOps, ImageFilter
    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False

try:
    _HAS_PYMUPDF = True  # fitz already imported above
except Exception:
    _HAS_PYMUPDF = False


# =========================================================
# P0：唯一 data URL helper（全檔案統一使用此函數）
# =========================================================
def bytes_to_data_url(b: bytes, mime: str) -> str:
    """
    回傳標準 data URL：data:{mime};base64,{b64}
    確保 mime 非空；若為空則 fallback 至 application/octet-stream。
    """
    if not mime:
        mime = "application/octet-stream"
    b64 = base64.b64encode(b).decode("utf-8")
    return f"data:{mime};base64,{b64}"


# =========================================================
# Text helpers
# =========================================================
def _clean_text(s: str) -> str:
    if not s:
        return ""
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    # 移除全角空格混亂問題
    s = re.sub(r"[\u3000]+", " ", s)
    return s.strip()


# P2：擴充可讀字符集，加入理科常用符號 + 常見中文標點
_READABLE_PATTERN = re.compile(
    r"[A-Za-z0-9\u4e00-\u9fff\u3400-\u4dbf"  # 基本漢字 + 擴展A區
    r"°μΩαβγδεζθλπσφψω"                     # 希臘字母
    r"\+\-×÷=<>\u2264\u2265\u2260\u00b1√∑∞∫∂∇"  # 數學符號（含 ≤ ≥ ≠ ±）
    r"→←↑↓⇌"                                # 化學/物理箭頭
    r"℃℉"                                   # 溫度單位
    r"\(\)\[\]{}|/\\^_~`"                   # 括號及常見符號
    r"，。！？；：""''（）【】《》、·…—￥"      # 中文標點
    r"｜「」『』〝〞〟※〰〾〿"                  # 中文特殊標點
    r"]"
)


def _text_quality_score(s: str) -> float:
    """
    P2：可讀字符比例 / 總長度
    （閾值建議 0.15，適合理科符號較多、OCR 碎片較多的情況）
    """
    if not s:
        return 0.0
    good = len(_READABLE_PATTERN.findall(s))
    return good / max(len(s), 1)


# P2：連續亂碼比例判斷
_GARBAGE_PATTERN = re.compile(
    r"[^\x09\x0a\x0d\x20-\x7e\u4e00-\u9fff\u3000-\u303f\uff00-\uffef\u3400-\u4dbf]{4,}"
)


def _is_garbage_text(s: str) -> bool:
    """True = OCR 結果太多亂碼，應丟棄。"""
    if not s:
        return True
    garbage_chars = sum(len(m) for m in _GARBAGE_PATTERN.findall(s))
    return (garbage_chars / max(len(s), 1)) > 0.40


# =========================================================
# P3：繁體中文 OCR 後處理（保守：預設不把中文標點英文化）
# =========================================================
def _postprocess_chinese_ocr(text: str, normalize_punct: bool = False) -> str:
    """
    修正常見繁體中文 OCR 錯誤：
    - 全角數字/字母轉半角（建議保留）
    -（可選）中文標點轉半角（預設 False，避免破壞中文語感）
    - 中文斷行合併：標點後的多餘換行移除
    """
    if not text:
        return ""

    # 全角英數轉半角
    fix_map = {
        "０": "0", "１": "1", "２": "2", "３": "3", "４": "4",
        "５": "5", "６": "6", "７": "7", "８": "8", "９": "9",
        "Ａ": "A", "Ｂ": "B", "Ｃ": "C", "Ｄ": "D", "Ｅ": "E",
        "Ｆ": "F", "Ｇ": "G", "Ｈ": "H", "Ｉ": "I", "Ｊ": "J",
        "Ｋ": "K", "Ｌ": "L", "Ｍ": "M", "Ｎ": "N", "Ｏ": "O",
        "Ｐ": "P", "Ｑ": "Q", "Ｒ": "R", "Ｓ": "S", "Ｔ": "T",
        "Ｕ": "U", "Ｖ": "V", "Ｗ": "W", "Ｘ": "X", "Ｙ": "Y",
        "Ｚ": "Z",
        "ａ": "a", "ｂ": "b", "ｃ": "c", "ｄ": "d", "ｅ": "e",
        "ｆ": "f", "ｇ": "g", "ｈ": "h", "ｉ": "i", "ｊ": "j",
        "ｋ": "k", "ｌ": "l", "ｍ": "m", "ｎ": "n", "ｏ": "o",
        "ｐ": "p", "ｑ": "q", "ｒ": "r", "ｓ": "s", "ｔ": "t",
        "ｕ": "u", "ｖ": "v", "ｗ": "w", "ｘ": "x", "ｙ": "y",
        "ｚ": "z",
    }

    if normalize_punct:
        fix_map.update({
            "，": ",", "。": ".", "！": "!", "？": "?", "；": ";",
            "：": ":", "（": "(", "）": ")", "【": "[", "】": "]",
            "—": "-", "～": "~", "｀": "`", "＂": "\"", "＇": "'",
        })

    for old, new in fix_map.items():
        text = text.replace(old, new)

    # 合併中文斷行（移除標點後的多餘換行）
    text = re.sub(r"([，。！？；：])\n", r"\1", text)
    return text


# =========================================================
# File text extractors
# =========================================================
def _extract_pdf_text(data: bytes) -> str:
    parts = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            parts.append(page.get_text())
    return _clean_text("\n".join(parts))


def _extract_docx_text(data: bytes) -> str:
    document = docx.Document(io.BytesIO(data))
    text = "\n".join(p.text for p in document.paragraphs if p.text)
    return _clean_text(text)


def _extract_txt_text(data: bytes) -> str:
    try:
        return _clean_text(data.decode("utf-8"))
    except Exception:
        return _clean_text(data.decode("utf-8", errors="ignore"))


def _extract_xlsx_text(data: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    chunks = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is None:
                    continue
                s = str(cell).strip()
                if s:
                    chunks.append(s)
    return _clean_text("\n".join(chunks))


def _extract_pptx_text(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                s = str(shape.text).strip()
                if s:
                    chunks.append(s)
    return _clean_text("\n".join(chunks))


# =========================================================
# OCR helpers（Streamlit Cloud / 中學掃描件向）
# =========================================================
def _safe_get_available_tess_langs() -> set:
    """
    嘗試查詢 tesseract 可用語言包；失敗就回傳空 set。
    （Streamlit Cloud / 容器環境有時會缺包或不可執行 list）
    """
    if not OCR_AVAILABLE:
        return set()
    try:
        langs = pytesseract.get_languages(config="")
        return set(langs or [])
    except Exception:
        return set()


def _normalize_ocr_lang(requested: str, fallback: str = "eng") -> str:
    """
    根據實際可用語言包自動降級：
    - requested: 例如 "chi_tra+chi_sim+eng"
    - 若某語言包不存在，移除
    - 最後若空，fallback 至 eng
    """
    req = (requested or "").strip()
    if not req:
        return fallback

    want = [x.strip() for x in req.split("+") if x.strip()]
    avail = _safe_get_available_tess_langs()
    if not avail:
        # 無法查詢可用語言：先原樣回傳，交由 tesseract 嘗試
        return req

    kept = [x for x in want if x in avail]
    if not kept:
        return fallback if fallback in avail else (want[-1] if want else fallback)
    return "+".join(kept)


def _otsu_threshold(gray_img: "Image.Image") -> int:
    """
    Otsu 二值化門檻（PIL histogram 版本，避免引入 OpenCV/Numpy）
    對掃描件底色不均比固定 threshold 更穩。
    """
    hist = gray_img.histogram()
    if not hist or len(hist) < 256:
        return 180

    total = sum(hist[:256])
    if total <= 0:
        return 180

    sum_total = 0
    for i in range(256):
        sum_total += i * hist[i]

    sum_b = 0
    w_b = 0
    max_var = -1.0
    threshold = 180

    for t in range(256):
        w_b += hist[t]
        if w_b == 0:
            continue
        w_f = total - w_b
        if w_f == 0:
            break
        sum_b += t * hist[t]
        m_b = sum_b / w_b
        m_f = (sum_total - sum_b) / w_f
        between = w_b * w_f * (m_b - m_f) ** 2
        if between > max_var:
            max_var = between
            threshold = t

    return int(threshold)


def _preprocess_image_for_ocr(
    img: "Image.Image",
    force_resize_2x: bool = True,
    binarize: bool = True,
) -> "Image.Image":
    """
    繁體中文掃描件向預處理：
    1) 灰階
    2) 中值濾波去噪
    3) autocontrast（cutoff=2）
    4) 二值化（預設用 Otsu 門檻）
    5) 2x 放大（提高小字辨識率）
    """
    # 1. 灰階
    img = img.convert("L")

    # 2. 去噪
    img = img.filter(ImageFilter.MedianFilter(size=3))

    # 3. 增強對比
    img = ImageOps.autocontrast(img, cutoff=2)

    # 4. 二值化（Otsu）
    if binarize:
        try:
            thr = _otsu_threshold(img)
        except Exception:
            thr = 180
        img = img.point(lambda x: 255 if x > thr else 0)

    # 5. 2倍縮放
    if force_resize_2x:
        w, h = img.size
        # 避免極大圖爆內存：限制最大邊長
        max_side = max(w, h)
        if max_side < 2400:
            img = img.resize((w * 2, h * 2), Image.Resampling.LANCZOS)

    return img


def _ocr_with_configs(
    img: "Image.Image",
    lang: str,
    psm_list=(6, 3),
    oem: int = 1,
    timeout_sec: int = 20,
) -> str:
    """
    以多個 PSM 進行 fallback 嘗試，返回最好（或第一個合格）的結果。
    """
    if not OCR_AVAILABLE:
        return ""

    lang = _normalize_ocr_lang(lang, fallback="eng")

    best_text = ""
    best_score = 0.0

    for psm in psm_list:
        try:
            custom_config = fr'--oem {int(oem)} --psm {int(psm)} -c preserve_interword_spaces=0'
            t = pytesseract.image_to_string(img, lang=lang, config=custom_config, timeout=timeout_sec)
            t = _clean_text(t)

            # 質量過濾
            if _text_quality_score(t) < 0.15 or _is_garbage_text(t):
                t_ok = ""
            else:
                t_ok = t

            score = _text_quality_score(t_ok) if t_ok else 0.0
            if score > best_score:
                best_score = score
                best_text = t_ok

            # 若已經有不錯文本，直接返回（節省時間）
            if best_text and best_score >= 0.22:
                return best_text
        except Exception:
            continue

    return best_text or ""


def _ocr_image_bytes(
    image_bytes: bytes,
    lang: str = "chi_tra+chi_sim+eng",
    normalize_punct: bool = False,
) -> str:
    """
    繁體中文 OCR 核心：
    - 預設語言：chi_tra + chi_sim + eng（會自動根據環境語言包降級）
    - PSM fallback：先 6，再 3
    - 後處理：全角英數轉半角、合併標點斷行（預設不英文化中文標點）
    """
    if not OCR_AVAILABLE:
        return ""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        img = _preprocess_image_for_ocr(img, force_resize_2x=True, binarize=True)

        text = _ocr_with_configs(
            img=img,
            lang=lang,
            psm_list=(6, 3),
            oem=1,
            timeout_sec=20,
        )
        if not text:
            return ""

        text = _postprocess_chinese_ocr(text, normalize_punct=normalize_punct)
        text = _clean_text(text)
        return text
    except Exception:
        return ""


# =========================================================
# Vision / data URL helpers（全部使用 bytes_to_data_url）
# =========================================================
def _pdf_pages_to_images_data_url(
    data: bytes,
    max_pages: int = 3,
    zoom: float = 2.0,
    dpi: int | None = None,
):
    """
    Vision 用：PDF 前 max_pages 頁渲染成 PNG data URL。
    優先使用 dpi（若 PyMuPDF 支援），否則 fallback Matrix(zoom, zoom)。
    """
    imgs = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        n = min(len(doc), max_pages)
        for i in range(n):
            page = doc[i]
            try:
                if dpi:
                    pix = page.get_pixmap(dpi=int(dpi), alpha=False)
                else:
                    mat = fitz.Matrix(float(zoom), float(zoom))
                    pix = page.get_pixmap(matrix=mat, alpha=False)
            except Exception:
                mat = fitz.Matrix(float(zoom), float(zoom))
                pix = page.get_pixmap(matrix=mat, alpha=False)

            imgs.append(bytes_to_data_url(pix.tobytes("png"), "image/png"))
    return imgs


def _pdf_page_to_png_bytes(
    page: "fitz.Page",
    dpi: int = 300,
    zoom: float = 2.5,
) -> bytes:
    """
    將 PDF page 渲染成 PNG bytes：
    優先用 dpi=300（OCR 常用），失敗再 fallback 用 Matrix zoom。
    """
    try:
        pix = page.get_pixmap(dpi=int(dpi), alpha=False)
    except Exception:
        mat = fitz.Matrix(float(zoom), float(zoom))
        pix = page.get_pixmap(matrix=mat, alpha=False)
    return pix.tobytes("png")


# =========================================================
# Main: extract_payload
# =========================================================
def extract_payload(
    file,
    enable_ocr: bool = False,
    ocr_lang: str = "chi_tra+chi_sim+eng",
    enable_vision: bool = False,
    vision_pdf_max_pages: int = 3,
    # P4（新增）：掃描 PDF OCR 控制參數（避免全本 OCR 拖死）
    ocr_pdf_max_pages: int = 10,
    ocr_pdf_dpi: int = 300,
    ocr_pdf_zoom_fallback: float = 2.5,
    ocr_min_text_len_for_skip: int = 50,
    normalize_punct: bool = False,
) -> dict:
    """
    回傳：{ "text": str, "images": [data_url,...], "meta": {...} }

    - enable_ocr：是否啟用 OCR
    - enable_vision：是否輸出 images（供多模態 LLM）
    - ocr_pdf_max_pages：掃描 PDF OCR 最多處理幾頁（避免 50 頁講義跑死）
    - ocr_pdf_dpi：PDF 轉圖 DPI（建議 300）
    """
    name = getattr(file, "name", "")
    ext = name.split(".")[-1].lower()
    data = file.getvalue()
    out = {"text": "", "images": [], "meta": {"ext": ext}}

    try:
        if ext == "pdf":
            text = _extract_pdf_text(data)
            out["text"] = text

            # 掃描 PDF：抽字太少 → OCR / Vision
            if len(text) < int(ocr_min_text_len_for_skip):
                if enable_ocr and OCR_AVAILABLE:
                    parts = []
                    with fitz.open(stream=data, filetype="pdf") as doc:
                        n_pages = len(doc)
                        max_pages = int(ocr_pdf_max_pages) if ocr_pdf_max_pages else n_pages
                        max_pages = max(1, min(n_pages, max_pages))

                        for i in range(max_pages):
                            page = doc[i]
                            png_bytes = _pdf_page_to_png_bytes(
                                page, dpi=int(ocr_pdf_dpi), zoom=float(ocr_pdf_zoom_fallback)
                            )
                            t = _ocr_image_bytes(png_bytes, lang=ocr_lang, normalize_punct=normalize_punct)
                            if t:
                                parts.append(t)

                    out["text"] = _clean_text("\n".join(parts)) if parts else ""

                if enable_vision:
                    try:
                        out["images"] = _pdf_pages_to_images_data_url(
                            data,
                            max_pages=int(vision_pdf_max_pages),
                            zoom=2.0,
                            dpi=None,
                        )
                    except Exception:
                        out["images"] = []

            return out

        if ext == "docx":
            out["text"] = _extract_docx_text(data)
            return out

        if ext == "txt":
            out["text"] = _extract_txt_text(data)
            return out

        if ext == "xlsx":
            out["text"] = _extract_xlsx_text(data)
            return out

        if ext == "pptx":
            out["text"] = _extract_pptx_text(data)
            return out

        # Image input
        if ext in {"png", "jpg", "jpeg"}:
            mime = "image/png" if ext == "png" else "image/jpeg"
            if enable_ocr:
                out["text"] = _ocr_image_bytes(data, lang=ocr_lang, normalize_punct=normalize_punct)
            if enable_vision:
                out["images"] = [bytes_to_data_url(data, mime)]
            return out

        return out
    except Exception:
        return out


# =========================================================
# P1：extract_images_for_llm_ocr 改呼叫共用 bytes_to_data_url
# =========================================================
def extract_images_for_llm_ocr(
    file,
    pdf_max_pages: int = 3,
    pdf_zoom: float = 2.0,
):
    """
    將圖片 / 掃描PDF（前N頁）轉成 data URL 供多模態 LLM 讀圖用。
    統一使用 bytes_to_data_url()，不再重複實作。
    """
    name = getattr(file, "name", "")
    ext = name.split(".")[-1].lower()
    data = file.getvalue()

    if ext in {"png", "jpg", "jpeg"}:
        mime = "image/png" if ext == "png" else "image/jpeg"
        return [bytes_to_data_url(data, mime)]

    if ext == "pdf":
        return _pdf_pages_to_images_data_url(
            data,
            max_pages=max(1, int(pdf_max_pages)),
            zoom=float(pdf_zoom),
            dpi=None,
        )

    return []


# =========================================================
# 兼容舊接口
# =========================================================
def extract_text(
    file,
    enable_ocr: bool = False,
    ocr_lang: str = "chi_tra+chi_sim+eng",
) -> str:
    return extract_payload(
        file,
        enable_ocr=enable_ocr,
        ocr_lang=ocr_lang,
        enable_vision=False,
    ).get("text", "")


# 兼容舊名稱
_image_bytes_to_data_url = bytes_to_data_url
_bytes_to_data_url = bytes_to_data_url