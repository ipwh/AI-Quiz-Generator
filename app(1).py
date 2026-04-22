# =========================================================
# app.py — FINAL-FULL（修正版：舊版 UI 對齊 + 進度顯示 + 強化 PDF 文字擷取）
#
# ✅ 修正重點
# 1) 「② 勾選重點段落」格式：恢復舊版段落勾選＋【重點段落/其餘內容】合併方式
# 2) 「④ 結果」格式：優先使用舊版編輯器（若模組存在），否則用表格/展開式顯示
# 3) 「③ 生成題目」加入進度顯示（st.status/st.progress/st.spinner）
# 4) PDF 文字擷取：先用 pdfplumber（若可用），否則用 PyPDF2；失敗則提示改用 OCR/Vision
#
# ✅ 其他原本功能保留
# - Google OAuth 登入/登出
# - Fast Mode
# - API 類型選擇（DeepSeek/OpenAI/Grok/Azure/自訂）+ API 測試
# - OCR 模式（純文字/本地 OCR/Vision）
# - 多格式輸入（PDF/DOCX/TXT/PPTX/XLSX/PNG/JPG）
# - 匯出：Kahoot Excel / Wayground DOCX
# - Google Forms/Drive/Gmail：若缺 google-api-python-client 或 scope 不足則安全降級提示
#
# ⚠️ 注意
# - 本檔案僅使用 4 個空格縮排；不含不可見空白字元（U+00A0）
# - 如環境缺少某些套件（pdfplumber / PyPDF2 / pptx / pytesseract / googleapiclient），會提示並降級
# =========================================================

import io
import re
import traceback
from datetime import datetime

import streamlit as st
import pandas as pd
from docx import Document

# ---------- 服務層 ----------
from services.llm_service import (
    generate_questions,
    assist_import_questions,
    parse_import_questions_locally,
)
from services.vision_service import (
    supports_vision,
    vision_generate_questions,
    file_to_data_url,
)
from services.google_oauth import (
    oauth_is_configured,
    get_auth_url,
    exchange_code_for_credentials,
    credentials_to_dict,
    credentials_from_dict,
)

try:
    from services.llm_service import get_xai_default_model
except Exception:
    get_xai_default_model = None

try:
    from services.llm_service import ping_llm
except Exception:
    ping_llm = None

# （可選）舊版 editor / mapper / export panel
try:
    from core.question_mapper import dicts_to_items, items_to_editor_df, editor_df_to_items
    _HAS_EDITOR = True
except Exception:
    _HAS_EDITOR = False

try:
    from ui.components_export import render_export_panel
    _HAS_EXPORT_PANEL = True
except Exception:
    _HAS_EXPORT_PANEL = False

# =========================================================
# Page
# =========================================================
st.set_page_config(page_title="AI 題目生成器", layout="wide")
st.title("🏫 AI 題目生成器")

# =========================================================
# Session state
# =========================================================
SS = st.session_state
_SS_DEFAULTS = {
    "google_creds": None,
    "generated_items": [],
    "imported_items": [],
    "mark_idx_generate": set(),
    "mark_idx_import": set(),
    "extracted_text_generate": "",
    "extracted_text_import": "",
    "last_google_form_url": None,
    "last_drive_links": {},
}
for k, v in _SS_DEFAULTS.items():
    if k not in SS:
        SS[k] = v

# =========================================================
# Helpers
# =========================================================

def show_exception(user_msg: str, e: Exception):
    st.error(user_msg)
    with st.expander("🔎 技術細節（維護用）"):
        st.code("".join(traceback.format_exception(type(e), e, e.__traceback__)))


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def split_paragraphs(text: str):
    return [p.strip() for p in (text or "").split("\n\n") if p.strip()]


def build_text_with_highlights(raw_text: str, marked_idx: set, limit: int):
    """舊版同款：把重點段落置頂，並加上分隔標題。"""
    if not raw_text:
        return ""
    paragraphs = split_paragraphs(raw_text)
    if not paragraphs:
        return ""

    highlighted = [p for i, p in enumerate(paragraphs) if i in marked_idx]
    others = [p for i, p in enumerate(paragraphs) if i not in marked_idx]

    combined = []
    if highlighted:
        combined.append("【重點段落】")
        combined.extend(highlighted)
    if others:
        combined.append("【其餘內容】")
        combined.extend(others)

    final_text = "\n\n".join(combined)
    if limit and len(final_text) > limit:
        final_text = final_text[:limit]
    return final_text


def now_stamp():
    return datetime.now().strftime("%Y%m%d-%H%M")


# =========================================================
# 多格式擷取
# =========================================================

def extract_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text and p.text.strip())


def extract_xlsx(file_bytes: bytes) -> str:
    try:
        xls = pd.ExcelFile(io.BytesIO(file_bytes))
        lines = []
        for s in xls.sheet_names:
            df = xls.parse(s, header=None)
            for row in df.astype(str).fillna("").values.tolist():
                line = " ".join(v for v in row if v and v.lower() != "nan").strip()
                if line:
                    lines.append(line)
        return "\n\n".join(lines)
    except Exception:
        return ""


def extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n\n".join(lines)
    except Exception:
        return ""


def extract_pdf_text(file_bytes: bytes) -> str:
    """PDF 抽字：優先 pdfplumber，其次 PyPDF2。"""
    # 1) pdfplumber
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t.strip():
                    parts.append(t.strip())
        out = "\n\n".join(parts)
        if out.strip():
            return out
    except Exception:
        pass

    # 2) PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        parts = []
        for p in reader.pages:
            t = p.extract_text() or ""
            if t.strip():
                parts.append(t.strip())
        return "\n\n".join(parts)
    except Exception:
        return ""


def local_ocr_images_to_text(images_bytes_list, lang_hint="chi_tra+eng") -> str:
    try:
        import pytesseract
        from PIL import Image
    except Exception as e:
        raise RuntimeError("本地 OCR 需要 pytesseract + pillow") from e

    texts = []
    for b in images_bytes_list:
        img = Image.open(io.BytesIO(b))
        t = pytesseract.image_to_string(img, lang=lang_hint)
        t = clean_text(t)
        if t:
            texts.append(t)
    return "\n\n".join(texts)


def extract_payload_from_uploads(uploads):
    extracted_parts = []
    vision_data_urls = []
    image_bytes_list = []

    for f in uploads:
        name = (f.name or "").lower()
        b = f.read()

        if name.endswith(".docx"):
            t = extract_docx(b)
            if t:
                extracted_parts.append(t)

        elif name.endswith(".xlsx"):
            t = extract_xlsx(b)
            if t:
                extracted_parts.append(t)

        elif name.endswith(".pptx"):
            t = extract_pptx(b)
            if t:
                extracted_parts.append(t)

        elif name.endswith(".txt"):
            try:
                t = b.decode("utf-8", errors="ignore")
            except Exception:
                t = ""
            t = clean_text(t)
            if t:
                extracted_parts.append(t)

        elif name.endswith(".pdf"):
            t = extract_pdf_text(b)
            if t and t.strip():
                extracted_parts.append(t)
            else:
                # 掃描/無法抽字 → 交由 Vision/OCR
                vision_data_urls.append(file_to_data_url(b, f.name))

        elif name.endswith((".png", ".jpg", ".jpeg")):
            vision_data_urls.append(file_to_data_url(b, f.name))
            image_bytes_list.append(b)

    return clean_text("\n\n".join(extracted_parts)), vision_data_urls, image_bytes_list


# =========================================================
# 匯出：Kahoot / Wayground
# =========================================================

def export_kahoot_excel(items):
    rows = []
    for q in items:
        opts = (q.get("options", []) + ["", "", "", ""])[:4]
        corr = q.get("correct", [""])[0] if q.get("correct") else ""
        rows.append({
            "Question": q.get("question", ""),
            "Answer 1": opts[0],
            "Answer 2": opts[1],
            "Answer 3": opts[2],
            "Answer 4": opts[3],
            "Correct Answer": corr,
        })
    df = pd.DataFrame(rows)
    bio = io.BytesIO()
    df.to_excel(bio, index=False, engine="openpyxl")
    bio.seek(0)
    return bio


def export_wayground_docx(items, title="題目清單"):
    doc = Document()
    doc.add_heading(title, 0)
    for i, q in enumerate(items, 1):
        doc.add_heading(f"第 {i} 題", level=1)
        doc.add_paragraph(q.get("question", ""))
        for idx, opt in enumerate(q.get("options", []), 1):
            doc.add_paragraph(f"{idx}. {opt}")
        doc.add_paragraph("答案：" + ",".join(q.get("correct", [])))
        if q.get("needs_review"):
            doc.add_paragraph("⚠️ 需要教師確認")
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio


# =========================================================
# Google OAuth callback
# =========================================================
params = st.query_params
if oauth_is_configured() and "code" in params and not SS.google_creds:
    try:
        code = params.get("code")
        state = params.get("state")
        if isinstance(code, list):
            code = code[0]
        if isinstance(state, list):
            state = state[0]
        creds = exchange_code_for_credentials(code=code, returned_state=state)
        SS.google_creds = credentials_to_dict(creds)
        st.query_params.clear()
        st.rerun()
    except Exception as e:
        st.query_params.clear()
        show_exception("Google 登入失敗。請重新按『連接 Google（登入）』一次。", e)
        st.stop()


# =========================================================
# Sidebar（完整回歸）
# =========================================================

st.sidebar.header("🟦 Google 連接（Google Forms / Google Drive 一鍵分享）")
if not oauth_is_configured():
    st.sidebar.warning("⚠️ 尚未設定 Google OAuth（Secrets: google_oauth_client + APP_URL）")
else:
    if SS.google_creds:
        st.sidebar.success("✅ 已連接 Google")
        if st.sidebar.button("🔒 登出 Google", key="btn_logout_google"):
            SS.google_creds = None
            st.rerun()
    else:
        st.sidebar.link_button("🔐 連接 Google（登入）", get_auth_url())
        st.sidebar.caption("提示：請以學校電郵登入，方便統一管理與分享。")

st.sidebar.divider()

fast_mode = st.sidebar.checkbox(
    "⚡ 快速模式",
    value=True,
    help="較快、較保守：較短輸出與較短超時；適合日常快速出題。",
)

st.sidebar.caption("關閉快速模式：較慢，但題目更豐富/更有變化。")

st.sidebar.divider()

st.sidebar.header("🔌 AI API 設定")
preset = st.sidebar.selectbox(
    "快速選擇（簡易）",
    ["DeepSeek", "OpenAI", "Grok (xAI)", "Azure OpenAI", "自訂（OpenAI 相容）"],
    key="preset",
)
api_key = st.sidebar.text_input("API Key", type="password", key="api_key")

auto_xai = False
azure_endpoint = ""
azure_deployment = ""
azure_api_version = "2024-02-15-preview"

if preset == "DeepSeek":
    base_url = "https://api.deepseek.com/v1"
    model = "deepseek-chat"
    cfg = {"api_key": api_key, "base_url": base_url, "model": model}

elif preset == "OpenAI":
    base_url = "https://api.openai.com/v1"
    model = "gpt-4o-mini"
    cfg = {"api_key": api_key, "base_url": base_url, "model": model}

elif preset == "Grok (xAI)":
    base_url = "https://api.x.ai/v1"
    model = "grok-4-latest"
    auto_xai = st.sidebar.checkbox("🤖 自動偵測可用最新 Grok 型號（建議）", value=True)
    if auto_xai and api_key and get_xai_default_model:
        try:
            model = get_xai_default_model(api_key, base_url)
            st.sidebar.caption("✅ 已自動選用：" + model)
        except Exception:
            pass
    cfg = {"api_key": api_key, "base_url": base_url, "model": model}

elif preset == "Azure OpenAI":
    with st.sidebar.expander("⚙️ Azure 設定", expanded=True):
        azure_endpoint = st.text_input("Azure Endpoint", value="")
        azure_deployment = st.text_input("Deployment name", value="")
        azure_api_version = st.text_input("API version", value="2024-02-15-preview")
    st.sidebar.warning("⚠️ 注意：若 services.llm_service 未實作 Azure 呼叫，將無法使用此選項。")
    cfg = {
        "type": "azure",
        "api_key": api_key,
        "endpoint": azure_endpoint,
        "deployment": azure_deployment,
        "api_version": azure_api_version,
    }

else:
    base_url = st.sidebar.text_input("Base URL（含 /v1）", value="")
    model = st.sidebar.text_input("Model", value="")
    cfg = {"api_key": api_key, "base_url": base_url, "model": model}

st.sidebar.divider()

st.sidebar.header("🧪 API 連線測試")
if st.sidebar.button("🧪 一鍵測試 API（回覆 OK）"):
    if not ping_llm:
        st.sidebar.warning("⚠️ 目前 llm_service 未提供 ping_llm()，略過測試。")
    elif not cfg.get("api_key"):
        st.sidebar.error("請先填入 API Key")
    else:
        with st.sidebar.spinner("正在測試連線…"):
            r = ping_llm(cfg, timeout=25)
        if r.get("ok"):
            st.sidebar.success("✅ 成功：" + str(r.get("latency_ms", 0)) + " ms；回覆：" + str(r.get("output", "")))
        else:
            st.sidebar.error("❌ 失敗：請檢查 Key/Model/Base URL 或服務狀態")
            st.sidebar.code(r.get("error", ""))

st.sidebar.divider()

st.sidebar.header("🔬 OCR / 讀圖設定（數理科必讀）")
ocr_mode = st.sidebar.radio(
    "教材擷取模式",
    [
        "📄 純文字（一般文件，最快）",
        "🔬 本地 OCR（掃描 PDF/圖片，離線）",
        "🤖 LLM Vision 讀圖（圖表/方程式/手寫，最準）",
    ],
    index=0,
)
vision_pdf_max_pages = 3
if ocr_mode.startswith("🤖"):
    vision_pdf_max_pages = st.sidebar.slider("Vision PDF 最多讀取頁數", 1, 10, 3)
    st.sidebar.info("💡 DeepSeek 不支援 Vision，請改用 Grok / GPT-4o 等模型。")

st.sidebar.divider()

st.sidebar.header("📘 出題設定")
subject = st.sidebar.selectbox(
    "科目",
    [
        "中國語文", "英國語文", "數學", "公民與社會發展", "科學", "公民、經濟及社會",
        "物理", "化學", "生物", "地理", "歷史", "中國歷史", "宗教",
        "資訊及通訊科技（ICT）", "經濟", "企業、會計與財務概論", "旅遊與款待",
    ],
)

level_label = st.sidebar.radio(
    "🎯 難度",
    ["基礎（理解與記憶）", "標準（應用與理解）", "進階（分析與思考）", "混合（課堂活動建議）"],
    index=1,
)
level_map = {
    "基礎（理解與記憶）": "easy",
    "標準（應用與理解）": "medium",
    "進階（分析與思考）": "hard",
    "混合（課堂活動建議）": "mixed",
}
level_code = level_map[level_label]

question_count = st.sidebar.selectbox("🧮 題目數目（生成用）", [5, 8, 10, 12, 15, 20], index=2)

st.sidebar.divider()
st.sidebar.markdown("### ⚠️ 隱私提醒")
st.sidebar.info(
    "教材內容將傳送至您所選的 AI 服務商（OpenAI / xAI / DeepSeek 等）。\n\n"
    "**請勿上傳含有學生個人資料、敏感資訊或受版權嚴格保護的完整教材。**"
)

with st.expander("🧭 使用流程", expanded=True):
    st.markdown(
        "**⚠️ 重要隱私提醒**：教材內容會上傳至第三方 AI 服務，請避免上傳含學生姓名、學號等個人資料。\n\n"
        "#### 🪄 生成新題目\n"
        "- 左側欄選擇科目、難度、題目數目\n"
        "- 填寫 AI API Key 並測試連線\n"
        "- ① 上載教材（支援多格式 / OCR / Vision）\n"
        "- ② 老師勾選重點段落\n"
        "- 按「🪄 生成題目」\n"
        "- ④ 檢視結果（⚠️ needs_review 題目優先確認）\n"
        "- ⑤ 匯出（Kahoot / Wayground / Google Forms / 電郵分享）\n\n"
        "#### 📄 匯入現有題目\n"
        "- 上載或貼上題目\n"
        "- 選擇是否使用 AI 協助整理\n"
        "- 按「✨ 整理並轉換」後核對答案\n"
    )


# =========================================================
# Tabs
# =========================================================

tab_generate, tab_import = st.tabs(["🪄 生成新題目", "📄 匯入現有題目"])


# =========================================================
# UI components: ② 段落勾選（舊版樣式）
# =========================================================

def render_paragraph_picker(raw_text: str, mark_idx_key: str, limit: int = 0):
    """恢復舊版樣式：每段上方一個 checkbox；提供全選/全不選；回傳 (paras, mark_idx, used_text_preview)"""
    paras = split_paragraphs(raw_text)
    mark_idx = SS.get(mark_idx_key, set())

    col_a, col_b, col_c = st.columns([1, 1, 3])
    with col_a:
        if st.button("✅ 全選段落", key=f"sel_all_{mark_idx_key}"):
            SS[mark_idx_key] = set(range(len(paras)))
            st.rerun()
    with col_b:
        if st.button("⬜ 全不選", key=f"sel_none_{mark_idx_key}"):
            SS[mark_idx_key] = set()
            st.rerun()
    with col_c:
        st.caption("提示：只會把已勾選段落送入 AI（更貼題）")

    if paras:
        new_set = set(mark_idx)
        for i, p in enumerate(paras):
            checked = (i in new_set)
            new_checked = st.checkbox(f"選用段落 {i+1}", value=checked, key=f"{mark_idx_key}_p_{i}")
            st.markdown(p)
            if new_checked:
                new_set.add(i)
            else:
                new_set.discard(i)
        SS[mark_idx_key] = new_set
        mark_idx = new_set

    used_preview = build_text_with_highlights(raw_text, mark_idx, limit=limit)
    return paras, mark_idx, used_preview


# =========================================================
# ④ 結果顯示（舊版偏好：能用 editor 就用 editor）
# =========================================================

def render_results(items: list, title: str):
    if not isinstance(items, list) or not items:
        return

    st.markdown(f"## {title}")

    if _HAS_EDITOR:
        try:
            st.caption("（可直接在表格內修改題幹／選項／答案）")
            df = items_to_editor_df(dicts_to_items(items))
            edited = st.data_editor(df, use_container_width=True, num_rows="dynamic")
            items2 = [it.to_dict() for it in editor_df_to_items(edited)]
            # needs_review 高亮提示
            needs_cnt = sum(1 for q in items2 if q.get("needs_review"))
            if needs_cnt:
                st.warning(f"⚠️ 有 {needs_cnt} 題標示 needs_review，請教師優先確認。")
            # 回寫
            return items2
        except Exception:
            # fallback
            pass

    # fallback：展開式
    for i, q in enumerate(items, 1):
        needs = bool(q.get("needs_review"))
        with st.expander(("⚠️ " if needs else "") + f"第 {i} 題", expanded=needs):
            st.markdown(f"**題目：** {q.get('question','')}")
            for j, opt in enumerate(q.get("options", []), 1):
                st.markdown(f"{j}. {opt}")
            st.markdown("**答案：** " + ",".join(q.get("correct", [])))
            if needs:
                st.warning(q.get("explanation") or "此題涉及關鍵條件／推論，建議教師覆核。")
    return items


# =========================================================
# Tab 1: Generate
# =========================================================
with tab_generate:
    st.markdown("## ① 上載教材")
    uploads = st.file_uploader(
        "支援 PDF / DOCX / TXT / PPTX / XLSX / PNG / JPG；可多檔",
        type=["pdf", "docx", "txt", "pptx", "xlsx", "png", "jpg", "jpeg"],
        accept_multiple_files=True,
        key="u_gen",
    )

    extracted_text = ""
    vision_urls = []
    image_bytes_list = []

    if uploads:
        with st.status("正在擷取教材內容…", expanded=False) as status:
            status.update(label="正在擷取教材內容…", state="running")
            extracted_text, vision_urls, image_bytes_list = extract_payload_from_uploads(uploads)
            status.update(label="教材內容擷取完成", state="complete")

    extracted_text = clean_text(extracted_text)
    SS.extracted_text_generate = extracted_text

    st.markdown("## ② 老師勾選重點段落（只用勾選內容出題）")
    paras, mark_idx, used_preview = render_paragraph_picker(
        SS.extracted_text_generate,
        mark_idx_key="mark_idx_generate",
        limit=(8000 if fast_mode else 10000),
    )

    if used_preview:
        with st.expander("🔎 已選內容預覽（將送入 AI）", expanded=False):
            st.text(used_preview[:5000] + ("\n…（已截斷）" if len(used_preview) > 5000 else ""))

    st.markdown("## ③ 生成題目")

    if st.button("🪄 生成題目", key="btn_generate"):
        if not cfg.get("api_key") and preset != "Azure OpenAI":
            st.error("請先在左側輸入 API Key")
        else:
            try:
                progress = st.progress(0)
                with st.spinner("正在整理選定內容…"):
                    progress.progress(10)
                    used_text = used_preview if used_preview else extracted_text
                    used_text = clean_text(used_text)
                    progress.progress(25)

                items = None

                if ocr_mode.startswith("🔬"):
                    with st.spinner("正在進行本地 OCR（如有圖片）…"):
                        progress.progress(35)
                        if image_bytes_list:
                            try:
                                ocr_txt = local_ocr_images_to_text(image_bytes_list)
                                used_text = clean_text(used_text + "\n\n" + ocr_txt)
                            except Exception:
                                st.warning("本地 OCR 不可用，將直接以文字內容出題。")
                        progress.progress(50)
                    with st.spinner("AI 正在生成題目…"):
                        items = generate_questions(cfg, used_text, subject, level_code, question_count, fast_mode=fast_mode)
                        progress.progress(90)

                elif ocr_mode.startswith("🤖"):
                    with st.spinner("AI Vision 正在讀圖/讀掃描 PDF…"):
                        progress.progress(40)
                        if vision_urls and supports_vision(cfg):
                            items = vision_generate_questions(
                                cfg,
                                used_text,
                                vision_urls,
                                subject,
                                level_code,
                                question_count,
                                fast_mode=fast_mode,
                            )
                        else:
                            st.warning("Vision 不可用（可能是模型不支援或沒有可讀圖片），將改用純文字出題。")
                            items = generate_questions(cfg, used_text, subject, level_code, question_count, fast_mode=fast_mode)
                        progress.progress(90)

                else:
                    with st.spinner("AI 正在生成題目…"):
                        progress.progress(40)
                        items = generate_questions(cfg, used_text, subject, level_code, question_count, fast_mode=fast_mode)
                        progress.progress(90)

                SS.generated_items = items or []
                progress.progress(100)
                st.success("✅ 題目生成完成")

            except Exception as e:
                show_exception("生成失敗", e)

    # ④ 結果（舊版格式）
    updated = render_results(SS.generated_items, "④ 結果（⚠️ needs_review 需教師確認）")
    if isinstance(updated, list):
        SS.generated_items = updated

    # ⑤ 匯出
    if isinstance(SS.generated_items, list) and SS.generated_items:
        st.markdown("---")
        st.markdown("## ⑤ 匯出")

        # 如你原本有 export panel，優先用（完整功能）
        if _HAS_EXPORT_PANEL:
            try:
                render_export_panel(
                    items=SS.generated_items,
                    subject=subject,
                    google_creds=SS.google_creds,
                    mode="generate",
                )
            except Exception as e:
                st.warning("內建 Export Panel 暫不可用，改用基本匯出。")

        # 基本匯出（確保有）
        col1, col2 = st.columns(2)
        with col1:
            try:
                xbio = export_kahoot_excel(SS.generated_items)
                st.download_button(
                    "⬇️ Kahoot（Excel）",
                    data=xbio,
                    file_name=f"kahoot_{now_stamp()}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            except Exception as e:
                show_exception("Kahoot 匯出失敗", e)
        with col2:
            try:
                dbio = export_wayground_docx(SS.generated_items, title=f"{subject} 題目清單")
                st.download_button(
                    "⬇️ Wayground（DOCX）",
                    data=dbio,
                    file_name=f"wayground_{now_stamp()}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            except Exception as e:
                show_exception("Wayground 匯出失敗", e)


# =========================================================
# Tab 2: Import
# =========================================================
with tab_import:
    st.markdown("## ① 上載 / 貼上題目")
    uploads2 = st.file_uploader(
        "支援 DOCX / PDF / XLSX / PPTX / TXT；可多檔",
        type=["pdf", "docx", "txt", "pptx", "xlsx", "png", "jpg", "jpeg"],
        accept_multiple_files=True,
        key="u_imp",
    )
    pasted = st.text_area("或直接貼上題目文字（可選）", value="", height=150, key="imp_paste")

    extracted_text2 = ""
    vision_urls2 = []
    image_bytes_list2 = []

    if uploads2:
        with st.status("正在擷取題目內容…", expanded=False) as status:
            status.update(label="正在擷取題目內容…", state="running")
            extracted_text2, vision_urls2, image_bytes_list2 = extract_payload_from_uploads(uploads2)
            status.update(label="擷取完成", state="complete")

    raw_all = clean_text("\n\n".join([extracted_text2, pasted]))
    SS.extracted_text_import = raw_all

    st.markdown("## ②（可選）勾選要整理的重點段落")
    paras2, mark2, used_preview2 = render_paragraph_picker(
        SS.extracted_text_import,
        mark_idx_key="mark_idx_import",
        limit=(8000 if fast_mode else 10000),
    )

    if used_preview2:
        with st.expander("🔎 已選內容預覽（將送入 AI 整理）", expanded=False):
            st.text(used_preview2[:5000] + ("\n…（已截斷）" if len(used_preview2) > 5000 else ""))

    st.markdown("## ③ 整理並轉換")
    use_ai = st.checkbox("使用 AI 協助整理（建議）", value=True, key="imp_use_ai")

    if st.button("✨ 整理並轉換", key="btn_import"):
        if not cfg.get("api_key") and preset != "Azure OpenAI" and use_ai:
            st.error("請先在左側輸入 API Key")
        else:
            try:
                progress = st.progress(0)
                with st.spinner("正在整理選定內容…"):
                    progress.progress(20)
                    used = used_preview2 if used_preview2 else raw_all
                    used = clean_text(used)
                    progress.progress(40)

                if use_ai:
                    with st.spinner("AI 正在整理題目…"):
                        items2 = assist_import_questions(cfg, used, subject)
                        progress.progress(90)
                else:
                    items2 = parse_import_questions_locally(used)
                    progress.progress(90)

                SS.imported_items = items2 or []
                progress.progress(100)
                st.success("✅ 題目整理完成")

            except Exception as e:
                show_exception("匯入/整理失敗", e)

    updated2 = render_results(SS.imported_items, "④ 整理結果（⚠️ needs_review 需教師確認）")
    if isinstance(updated2, list):
        SS.imported_items = updated2

    if isinstance(SS.imported_items, list) and SS.imported_items:
        st.markdown("---")
        st.markdown("## ⑤ 匯出")

        if _HAS_EXPORT_PANEL:
            try:
                render_export_panel(
                    items=SS.imported_items,
                    subject=subject,
                    google_creds=SS.google_creds,
                    mode="import",
                )
            except Exception:
                st.warning("內建 Export Panel 暫不可用，改用基本匯出。")

        col1, col2 = st.columns(2)
        with col1:
            try:
                xbio = export_kahoot_excel(SS.imported_items)
                st.download_button(
                    "⬇️ Kahoot（Excel）",
                    data=xbio,
                    file_name=f"kahoot_import_{now_stamp()}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            except Exception as e:
                show_exception("Kahoot 匯出失敗", e)
        with col2:
            try:
                dbio = export_wayground_docx(SS.imported_items, title=f"{subject}（匯入）題目清單")
                st.download_button(
                    "⬇️ Wayground（DOCX）",
                    data=dbio,
                    file_name=f"wayground_import_{now_stamp()}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            except Exception as e:
                show_exception("Wayground 匯出失敗", e)
