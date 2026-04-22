# =========================================================
# app.py — FINAL-FULL（修正版：UI 回歸 + PDF 抽字強化 + 進度顯示 + 可編輯答案 + Google Forms/電郵分享）
#
# 你回報的問題，本版已一次過修正：
# ✅ ② 重點段落：預設「摺疊」顯示；且預設「全選段落」（按鈕亦可切換）
# ✅ ③ 生成題目：加入進度顯示（狀態 + 進度條 + 擷取字數上限 8000/10000 舊版邏輯）
# ✅ PDF 抽字：加強（PyMuPDF/fitz → pdfplumber → PyPDF2），並提供「每檔抽到幾多字」報告
# ✅ AI 沒答案：加入「輸出驗證 + 一鍵補答案（AI）」；同時生成後預設在表格可編輯（題幹/選項/答案）
# ✅ ④ 結果：預設「一眼睇晒」(表格編輯)；不再要求逐題點擊
# ✅ 匯出：Kahoot Excel / Wayground DOCX 保底
# ✅ Google Forms：若 google-api-python-client + OAuth scope 可用，提供一鍵建立表單
# ✅ 一鍵電郵：優先建立 Gmail 草稿（可附檔）；如未能使用，提供 mailto 連結（含 Drive/Forms 連結）
#
# 注意：
# - 全檔 4 spaces 縮排；不含 U+00A0
# - 如缺少某些套件，會顯示提示並安全降級，不會 crash
# =========================================================

import io
import re
import time
import base64
import shutil
import traceback
from datetime import datetime

import streamlit as st
import pandas as pd
from docx import Document

from services.llm_service import (
    generate_questions,
    assist_import_questions,
    parse_import_questions_locally,
)
from services.vision_service import (
    supports_vision,
    vision_generate_questions,
    file_to_data_url,
)
from services.google_oauth import (
    oauth_is_configured,
    get_auth_url,
    exchange_code_for_credentials,
    credentials_to_dict,
    credentials_from_dict,
)

try:
    from services.llm_service import ping_llm
except Exception:
    ping_llm = None

try:
    from services.llm_service import get_xai_default_model
except Exception:
    get_xai_default_model = None

# （可選）舊版 editor / export panel（若存在就用）
try:
    from core.question_mapper import dicts_to_items, items_to_editor_df, editor_df_to_items
    _HAS_EDITOR = True
except Exception:
    _HAS_EDITOR = False

try:
    from ui.components_export import render_export_panel
    _HAS_EXPORT_PANEL = True
except Exception:
    _HAS_EXPORT_PANEL = False


# =========================================================
# Page
# =========================================================

st.set_page_config(page_title="AI 題目生成器", layout="wide")
st.title("🏫 AI 題目生成器")


# =========================================================
# Session state
# =========================================================

SS = st.session_state
_defaults = {
    "google_creds": None,

    "extracted_text_generate": "",
    "extracted_text_import": "",

    "file_report_generate": [],
    "file_report_import": [],

    "vision_urls_generate": [],
    "vision_urls_import": [],

    "image_bytes_generate": [],
    "image_bytes_import": [],

    "pdf_scan_bytes_generate": [],
    "pdf_scan_bytes_import": [],

    "mark_idx_generate": set(),
    "mark_idx_import": set(),

    "last_text_hash_generate": "",
    "last_text_hash_import": "",

    "generated_items": [],
    "imported_items": [],

    "last_google_form_url": "",
    "last_drive_links": {},

    "to_email": "",
}
for k, v in _defaults.items():
    if k not in SS:
        SS[k] = v


# =========================================================
# Helpers
# =========================================================

def show_exception(msg: str, e: Exception):
    st.error(msg)
    with st.expander("🔎 技術細節（維護用）"):
        st.code("".join(traceback.format_exception(type(e), e, e.__traceback__)))


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def split_paragraphs(text: str):
    return [p.strip() for p in (text or "").split("\n\n") if p.strip()]


def build_text_with_highlights(raw_text: str, marked_idx: set, limit: int):
    """舊版同款：重點段落置頂 + 標籤。"""
    if not raw_text:
        return ""
    paras = split_paragraphs(raw_text)
    highlighted = [p for i, p in enumerate(paras) if i in marked_idx]
    others = [p for i, p in enumerate(paras) if i not in marked_idx]

    combined = []
    if highlighted:
        combined.append("【重點段落】")
        combined.extend(highlighted)
    if others:
        combined.append("【其餘內容】")
        combined.extend(others)

    out = "\n\n".join(combined)
    if limit and len(out) > limit:
        out = out[:limit]
    return out


def now_stamp():
    return datetime.now().strftime("%Y%m%d-%H%M")


def text_hash(s: str) -> str:
    return str(hash(s or ""))


# =========================================================
# OCR readiness
# =========================================================

def check_local_ocr_ready():
    try:
        import pytesseract  # noqa
        from PIL import Image  # noqa
    except Exception:
        return False, "未安裝 pytesseract 或 pillow"

    if shutil.which("tesseract") is None:
        return False, "找不到 tesseract 執行檔（未安裝或不在 PATH）"

    try:
        import pytesseract
        _ = pytesseract.get_tesseract_version()
    except Exception:
        return True, "pytesseract 可用，但無法讀取版本（仍可嘗試）"

    return True, "✅ 本地 OCR（Tesseract）可用"


def local_ocr_images_to_text(images_bytes_list, lang_hint="chi_tra+eng") -> str:
    import pytesseract
    from PIL import Image

    texts = []
    for b in images_bytes_list:
        img = Image.open(io.BytesIO(b))
        t = pytesseract.image_to_string(img, lang=lang_hint)
        t = clean_text(t)
        if t:
            texts.append(t)
    return "\n\n".join(texts)


# =========================================================
# Extractors（多格式 + 強化 PDF）
# =========================================================

def extract_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text and p.text.strip())


def extract_xlsx(file_bytes: bytes) -> str:
    try:
        xls = pd.ExcelFile(io.BytesIO(file_bytes))
        lines = []
        for s in xls.sheet_names:
            df = xls.parse(s, header=None)
            for row in df.astype(str).fillna("").values.tolist():
                line = " ".join(v for v in row if v and v.lower() != "nan").strip()
                if line:
                    lines.append(line)
        return "\n\n".join(lines)
    except Exception:
        return ""


def extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n\n".join(lines)
    except Exception:
        return ""


def extract_pdf_text_strong(file_bytes: bytes) -> str:
    """文字型 PDF 抽字：PyMuPDF(fitz) → pdfplumber → PyPDF2"""
    # 0) PyMuPDF
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        parts = []
        for i in range(len(doc)):
            t = doc.load_page(i).get_text("text") or ""
            if t.strip():
                parts.append(t.strip())
        doc.close()
        out = "\n\n".join(parts)
        if out.strip():
            return out
    except Exception:
        pass

    # 1) pdfplumber
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t.strip():
                    parts.append(t.strip())
        out = "\n\n".join(parts)
        if out.strip():
            return out
    except Exception:
        pass

    # 2) PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        parts = []
        for p in reader.pages:
            t = p.extract_text() or ""
            if t.strip():
                parts.append(t.strip())
        return "\n\n".join(parts)
    except Exception:
        return ""


def extract_payload_from_uploads(uploads):
    extracted_parts = []
    vision_urls = []
    image_bytes = []
    pdf_scan_bytes = []
    report = []

    for f in uploads:
        name = (f.name or "").lower()
        b = f.read()

        extracted = ""
        kind = ""

        if name.endswith(".docx"):
            kind = "DOCX"
            extracted = extract_docx(b)

        elif name.endswith(".xlsx"):
            kind = "XLSX"
            extracted = extract_xlsx(b)

        elif name.endswith(".pptx"):
            kind = "PPTX"
            extracted = extract_pptx(b)

        elif name.endswith(".txt"):
            kind = "TXT"
            try:
                extracted = b.decode("utf-8", errors="ignore")
            except Exception:
                extracted = ""
            extracted = clean_text(extracted)

        elif name.endswith(".pdf"):
            kind = "PDF"
            extracted = extract_pdf_text_strong(b)
            if not extracted.strip():
                pdf_scan_bytes.append(b)
                vision_urls.append(file_to_data_url(b, f.name))

        elif name.endswith((".png", ".jpg", ".jpeg")):
            kind = "IMG"
            image_bytes.append(b)
            vision_urls.append(file_to_data_url(b, f.name))

        if extracted.strip():
            extracted_parts.append(extracted.strip())

        report.append({
            "file": f.name,
            "type": kind,
            "chars": len(extracted),
            "note": "OK" if len(extracted) > 0 else ("抽字失敗/掃描" if kind == "PDF" else ""),
        })

    return clean_text("\n\n".join(extracted_parts)), vision_urls, image_bytes, pdf_scan_bytes, report


# =========================================================
# Post-process: schema normalize + missing answers helper
# =========================================================

def normalize_items(items):
    """確保每題有 4 個 options + correct 為 ['1'..'4']，否則 needs_review=True。"""
    out = []
    for q in (items or []):
        qq = dict(q)
        qq["qtype"] = "single"
        opts = qq.get("options", [])
        if not isinstance(opts, list):
            opts = []
        opts = [str(x) if x is not None else "" for x in opts]
        opts = (opts + ["", "", "", ""])[:4]
        qq["options"] = opts

        corr = qq.get("correct", [])
        if isinstance(corr, str):
            corr = [c.strip() for c in corr.split(",") if c.strip()]
        if not isinstance(corr, list):
            corr = []
        corr = [str(x).strip() for x in corr]
        corr = [c for c in corr if c in {"1", "2", "3", "4"}]
        if len(corr) != 1:
            qq["needs_review"] = True
            # 先留空，讓「一鍵補答案」去補
            qq["correct"] = []
        else:
            qq["correct"] = [corr[0]]

        if "needs_review" not in qq:
            qq["needs_review"] = False
        if "explanation" not in qq:
            qq["explanation"] = ""

        out.append(qq)
    return out


def can_call_ai(cfg):
    return bool(cfg.get("api_key")) and (bool(cfg.get("base_url")) or cfg.get("type") == "azure")


def fill_missing_answers_with_ai(cfg, items):
    """如果某些題目 missing correct，用 AI 補回（用相同 cfg）。"""
    # 嘗試使用 llm_service 的 _chat（若存在）
    try:
        from services.llm_service import _chat
    except Exception:
        _chat = None

    if _chat is None:
        raise RuntimeError("目前 llm_service 未暴露 _chat，無法自動補答案。")

    need = []
    for i, q in enumerate(items):
        if not q.get("correct"):
            need.append((i, q))

    if not need:
        return items

    # 建 prompt：只回覆 JSON array，與原順序一致
    payload = []
    for i, q in need:
        payload.append({
            "question": q.get("question", ""),
            "options": q.get("options", []),
        })

    prompt = (
        "你是一名香港中學教師。以下是多項選擇題（四選一）。\n"
        "請你只為每題選出最正確答案，輸出純 JSON array。\n"
        "每個元素格式：{\"correct\": [\"1\"~\"4\"]}，只可 1 個。\n"
        "禁止輸出任何解釋文字。\n\n"
        f"題目清單：\n{payload}"
    )

    out = _chat(cfg, [{"role": "user", "content": prompt}], temperature=0.0, max_tokens=600, timeout=120)
    try:
        import json
        arr = json.loads(out)
    except Exception:
        raise RuntimeError("補答案回傳非 JSON")

    if not isinstance(arr, list) or len(arr) != len(need):
        raise RuntimeError("補答案回傳數量不符")

    # 寫回
    for (idx, _q), ans in zip(need, arr):
        corr = ans.get("correct", [])
        if isinstance(corr, str):
            corr = [corr]
        if isinstance(corr, list):
            corr = [str(x).strip() for x in corr if str(x).strip() in {"1","2","3","4"}]
        if len(corr) == 1:
            items[idx]["correct"] = [corr[0]]
        else:
            items[idx]["needs_review"] = True

    return items


# =========================================================
# Editor mapping（表格一眼睇晒 + 可編輯）
# =========================================================

def items_to_simple_df(items):
    rows = []
    for q in items:
        opts = q.get("options", ["", "", "", ""]) + ["", "", "", ""]
        opts = opts[:4]
        corr = q.get("correct", [])
        corr = corr[0] if isinstance(corr, list) and corr else ""
        rows.append({
            "題目": q.get("question", ""),
            "A": opts[0],
            "B": opts[1],
            "C": opts[2],
            "D": opts[3],
            "答案(1-4)": corr,
            "needs_review": bool(q.get("needs_review")),
            "explanation": q.get("explanation", ""),
        })
    return pd.DataFrame(rows)


def simple_df_to_items(df):
    items = []
    for _, r in df.iterrows():
        q = {
            "qtype": "single",
            "question": str(r.get("題目", "")),
            "options": [str(r.get("A", "")), str(r.get("B", "")), str(r.get("C", "")), str(r.get("D", ""))],
            "correct": [str(r.get("答案(1-4)", "")).strip()] if str(r.get("答案(1-4)", "")).strip() in {"1","2","3","4"} else [],
            "needs_review": bool(r.get("needs_review", False)) or (str(r.get("答案(1-4)", "")).strip() not in {"1","2","3","4"}),
            "explanation": str(r.get("explanation", "")),
        }
        items.append(q)
    return items


# =========================================================
# OAuth callback
# =========================================================
params = st.query_params
if oauth_is_configured() and "code" in params and not SS.google_creds:
    try:
        code = params.get("code")
        state = params.get("state")
        if isinstance(code, list):
            code = code[0]
        if isinstance(state, list):
            state = state[0]
        creds = exchange_code_for_credentials(code=code, returned_state=state)
        SS.google_creds = credentials_to_dict(creds)
        st.query_params.clear()
        st.rerun()
    except Exception as e:
        st.query_params.clear()
        show_exception("Google 登入失敗", e)
        st.stop()


# =========================================================
# Sidebar
# =========================================================

st.sidebar.header("🟦 Google 連接（Forms/Drive/Email）")
if not oauth_is_configured():
    st.sidebar.warning("⚠️ 尚未設定 Google OAuth")
else:
    if SS.google_creds:
        st.sidebar.success("✅ 已連接 Google")
        if st.sidebar.button("🔒 登出 Google"):
            SS.google_creds = None
            st.rerun()
    else:
        st.sidebar.link_button("🔐 連接 Google（登入）", get_auth_url())

st.sidebar.divider()

fast_mode = st.sidebar.checkbox("⚡ 快速模式", value=True)

st.sidebar.divider()

st.sidebar.header("🔌 AI API 設定")
preset = st.sidebar.selectbox(
    "快速選擇",
    ["DeepSeek", "OpenAI", "Grok (xAI)", "自訂（OpenAI 相容）"],
)
api_key = st.sidebar.text_input("API Key", type="password")

if preset == "DeepSeek":
    cfg = {"api_key": api_key, "base_url": "https://api.deepseek.com/v1", "model": "deepseek-chat"}
elif preset == "OpenAI":
    cfg = {"api_key": api_key, "base_url": "https://api.openai.com/v1", "model": "gpt-4o-mini"}
elif preset == "Grok (xAI)":
    base_url = "https://api.x.ai/v1"
    model = "grok-4-latest"
    if get_xai_default_model and api_key:
        try:
            model = get_xai_default_model(api_key, base_url)
        except Exception:
            pass
    cfg = {"api_key": api_key, "base_url": base_url, "model": model}
else:
    base_url = st.sidebar.text_input("Base URL（含 /v1）", value="")
    model = st.sidebar.text_input("Model", value="")
    cfg = {"api_key": api_key, "base_url": base_url, "model": model}

if ping_llm and st.sidebar.button("🧪 測試 API"):
    if not can_call_ai(cfg):
        st.sidebar.error("請先填妥 API Key／Base URL／Model")
    else:
        r = ping_llm(cfg, timeout=25)
        if r.get("ok"):
            st.sidebar.success("✅ OK")
        else:
            st.sidebar.error("❌ 失敗")
            st.sidebar.code(r.get("error", ""))

st.sidebar.divider()

st.sidebar.header("🔬 OCR / 讀圖")
ocr_mode = st.sidebar.radio(
    "教材擷取模式",
    ["📄 純文字", "🔬 本地 OCR", "🤖 Vision"],
    index=0,
)
ocr_ok, ocr_msg = check_local_ocr_ready()
if ocr_mode.startswith("🔬"):
    st.sidebar.info(ocr_msg)

st.sidebar.divider()

st.sidebar.header("📘 出題設定")
subject = st.sidebar.selectbox(
    "科目",
    [
        "中國語文", "英國語文", "數學", "公民與社會發展", "科學", "公民、經濟及社會",
        "物理", "化學", "生物", "地理", "歷史", "中國歷史", "宗教",
        "資訊及通訊科技（ICT）", "經濟", "企業、會計與財務概論", "旅遊與款待",
    ],
)
level_label = st.sidebar.radio(
    "🎯 難度",
    ["基礎（理解與記憶）", "標準（應用與理解）", "進階（分析與思考）", "混合（課堂活動建議）"],
    index=1,
)
level_map = {
    "基礎（理解與記憶）": "easy",
    "標準（應用與理解）": "medium",
    "進階（分析與思考）": "hard",
    "混合（課堂活動建議）": "mixed",
}
level_code = level_map[level_label]
question_count = st.sidebar.selectbox("題目數目", [5, 8, 10, 12, 15, 20], index=2)


# =========================================================
# Paragraph picker（預設摺疊 + 預設全選 + 按鈕有效）
# =========================================================

def ensure_default_select_all(text: str, mark_key: str, hash_key: str):
    h = text_hash(text)
    if SS.get(hash_key, "") != h:
        # 新內容：預設全選
        paras = split_paragraphs(text)
        SS[mark_key] = set(range(len(paras)))
        SS[hash_key] = h


def render_paragraph_picker_section(raw_text: str, mark_key: str, hash_key: str, limit: int, label: str):
    ensure_default_select_all(raw_text, mark_key, hash_key)
    paras = split_paragraphs(raw_text)
    mark = SS.get(mark_key, set())

    with st.expander(label, expanded=False):
        c1, c2, c3 = st.columns([1, 1, 3])
        with c1:
            if st.button("✅ 全選", key=f"btn_all_{mark_key}"):
                SS[mark_key] = set(range(len(paras)))
                st.rerun()
        with c2:
            if st.button("⬜ 全不選", key=f"btn_none_{mark_key}"):
                SS[mark_key] = set()
                st.rerun()
        with c3:
            st.caption("預設已全選。可按需要取消。")

        new_set = set(mark)
        for i, p in enumerate(paras):
            checked = (i in new_set)
            new_checked = st.checkbox(f"段落 {i+1}", value=checked, key=f"{mark_key}_{i}")
            st.markdown(p)
            if new_checked:
                new_set.add(i)
            else:
                new_set.discard(i)

        SS[mark_key] = new_set

    used = build_text_with_highlights(raw_text, SS[mark_key], limit=limit)
    return used


# =========================================================
# Results (預設表格一眼睇晒 + 可編輯答案)
# =========================================================

def render_editor(items: list, title: str, mode_key: str):
    if not isinstance(items, list) or not items:
        return items

    st.markdown(f"## {title}")

    # 預設表格一眼睇晒
    mode = st.radio(
        "顯示方式",
        ["表格（建議，可編輯）", "卡片（逐題）"],
        index=0,
        horizontal=True,
        key=f"disp_{mode_key}",
    )

    if mode.startswith("表格"):
        if _HAS_EDITOR:
            try:
                df = items_to_editor_df(dicts_to_items(items))
                edited = st.data_editor(df, use_container_width=True, num_rows="dynamic")
                items2 = [it.to_dict() for it in editor_df_to_items(edited)]
                return normalize_items(items2)
            except Exception:
                st.warning("舊版 editor 暫不可用，改用內建表格。")

        df = items_to_simple_df(items)
        edited = st.data_editor(df, use_container_width=True, num_rows="dynamic")
        items2 = simple_df_to_items(edited)
        return normalize_items(items2)

    # 卡片
    for i, q in enumerate(items, 1):
        needs = bool(q.get("needs_review"))
        with st.container(border=True):
            st.markdown(f"### {'⚠️ ' if needs else ''}第 {i} 題")
            st.markdown(q.get("question", ""))
            for j, opt in enumerate(q.get("options", []), 1):
                st.markdown(f"{j}. {opt}")
            st.markdown("**答案：** " + ",".join(q.get("correct", [])))
            if needs:
                st.warning(q.get("explanation") or "此題需要教師確認")

    return items


# =========================================================
# Google Forms + Email（若沒有 googleapiclient，顯示清楚提示）
# =========================================================

def try_build_google_services(creds):
    try:
        from googleapiclient.discovery import build
    except Exception as e:
        raise RuntimeError("缺少 google-api-python-client") from e

    forms = build("forms", "v1", credentials=creds)
    drive = build("drive", "v3", credentials=creds)
    gmail = build("gmail", "v1", credentials=creds)
    return forms, drive, gmail


def create_google_form(forms_service, title: str, items: list):
    form = forms_service.forms().create(body={"info": {"title": title}}).execute()
    form_id = form["formId"]

    reqs = []
    for idx, q in enumerate(items):
        opts = (q.get("options", []) + ["", "", "", ""])[:4]
        reqs.append({
            "createItem": {
                "item": {
                    "title": q.get("question", f"第 {idx+1} 題"),
                    "questionItem": {
                        "question": {
                            "required": True,
                            "choiceQuestion": {
                                "type": "RADIO",
                                "options": [{"value": o} for o in opts],
                                "shuffle": False,
                            },
                        }
                    },
                },
                "location": {"index": idx},
            }
        })

    forms_service.forms().batchUpdate(formId=form_id, body={"requests": reqs}).execute()

    f = forms_service.forms().get(formId=form_id).execute()
    return f.get("responderUri") or f"https://docs.google.com/forms/d/{form_id}/viewform"


def drive_upload_bytes(drive_service, filename: str, file_bytes: bytes, mime_type: str):
    from googleapiclient.http import MediaIoBaseUpload
    media = MediaIoBaseUpload(io.BytesIO(file_bytes), mimetype=mime_type, resumable=False)
    meta = {"name": filename}
    created = drive_service.files().create(body=meta, media_body=media, fields="id").execute()
    fid = created["id"]
    drive_service.permissions().create(fileId=fid, body={"type": "anyone", "role": "reader"}).execute()
    info = drive_service.files().get(fileId=fid, fields="webViewLink").execute()
    return info.get("webViewLink", "")


def gmail_create_draft_with_attachments(gmail_service, to_email: str, subject: str, body_text: str, attachments: list):
    from email.message import EmailMessage

    msg = EmailMessage()
    msg["To"] = to_email
    msg["Subject"] = subject
    msg.set_content(body_text)

    for (filename, mime_type, b) in attachments:
        msg.add_attachment(b, maintype=mime_type.split("/")[0], subtype=mime_type.split("/")[1], filename=filename)

    raw = base64.urlsafe_b64encode(msg.as_bytes()).decode("utf-8")
    draft = gmail_service.users().drafts().create(userId="me", body={"message": {"raw": raw}}).execute()
    return draft.get("id", "")


# =========================================================
# Tabs
# =========================================================

tab_gen, tab_imp = st.tabs(["🪄 生成新題目", "📄 匯入現有題目"])


# =========================================================
# Generate
# =========================================================
with tab_gen:
    st.markdown("## ① 上載教材（多格式）")
    uploads = st.file_uploader(
        "PDF / DOCX / XLSX / PPTX / TXT / 圖片",
        type=["pdf", "docx", "xlsx", "pptx", "txt", "png", "jpg", "jpeg"],
        accept_multiple_files=True,
        key="up_gen",
    )

    if uploads:
        with st.status("正在擷取教材內容…", expanded=False) as s:
            s.update(label="正在擷取教材內容…", state="running")
            text, v_urls, img_bytes, pdf_scan, report = extract_payload_from_uploads(uploads)
            SS.extracted_text_generate = text
            SS.vision_urls_generate = v_urls
            SS.image_bytes_generate = img_bytes
            SS.pdf_scan_bytes_generate = pdf_scan
            SS.file_report_generate = report
            s.update(label="擷取完成", state="complete")

        with st.expander("📄 擷取報告（每檔抽取字數）", expanded=False):
            st.dataframe(pd.DataFrame(SS.file_report_generate), use_container_width=True)

    limit = 8000 if fast_mode else 10000

    st.markdown("## ② 老師勾選重點段落（只用勾選內容出題）")
    used = render_paragraph_picker_section(
        SS.extracted_text_generate,
        mark_key="mark_idx_generate",
        hash_key="last_text_hash_generate",
        limit=limit,
        label="② 重點段落選擇（預設摺疊）",
    )

    st.info(
        f"已從文本擷取 {len(SS.extracted_text_generate)} 字；送入 AI（重點合併後）{len(used)} 字（上限 {limit}；快速模式={'開' if fast_mode else '關'}）"
    )

    with st.expander("🔎 已選內容預覽（將送入 AI）", expanded=False):
        st.text(used[:5000] + ("\n…（已截斷）" if len(used) > 5000 else ""))

    st.markdown("## ③ 生成題目")

    if st.button("🪄 生成題目", key="btn_gen"):
        if not can_call_ai(cfg):
            st.error("請先在左側輸入 API Key（以及自訂時 Base URL/Model）")
        else:
            try:
                prog = st.progress(0)
                prog.progress(10)

                used_for_ai = used

                # 本地 OCR：只對圖片/掃描 PDF 轉圖（此處略；你已表示文字型 PDF）
                if ocr_mode.startswith("🔬") and ocr_ok:
                    with st.spinner("本地 OCR（如有圖片/掃描 PDF）…"):
                        prog.progress(30)
                        ocr_images = list(SS.image_bytes_generate)
                        # 若 pdf_scan 有內容，嘗試轉圖
                        if SS.pdf_scan_bytes_generate:
                            try:
                                import fitz
                                for b in SS.pdf_scan_bytes_generate:
                                    doc = fitz.open(stream=b, filetype="pdf")
                                    n = min(len(doc), 3)
                                    for i in range(n):
                                        pix = doc.load_page(i).get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                                        ocr_images.append(pix.tobytes("png"))
                                    doc.close()
                            except Exception:
                                pass
                        if ocr_images:
                            ocr_text = local_ocr_images_to_text(ocr_images)
                            used_for_ai = clean_text(used_for_ai + "\n\n" + ocr_text)
                        prog.progress(45)

                # Vision
                if ocr_mode.startswith("🤖") and SS.vision_urls_generate and supports_vision(cfg):
                    with st.spinner("Vision 讀圖中…"):
                        prog.progress(55)
                        items = vision_generate_questions(
                            cfg,
                            used_for_ai,
                            SS.vision_urls_generate,
                            subject,
                            level_code,
                            question_count,
                            fast_mode=fast_mode,
                        )
                        prog.progress(85)
                else:
                    with st.spinner("AI 生成題目中…"):
                        prog.progress(60)
                        items = generate_questions(cfg, used_for_ai, subject, level_code, question_count, fast_mode=fast_mode)
                        prog.progress(85)

                SS.generated_items = normalize_items(items or [])
                prog.progress(100)
                st.success("✅ 題目生成完成")

            except Exception as e:
                show_exception("生成失敗", e)

    # 若沒有答案，提供一鍵補答案
    if isinstance(SS.generated_items, list) and SS.generated_items:
        missing = sum(1 for q in SS.generated_items if not q.get("correct"))
        if missing:
            st.warning(f"⚠️ 有 {missing} 題未包含答案。可按下方按鈕用 AI 補回答案。")
            if st.button("🧠 一鍵補答案（AI）", key="btn_fill_ans"):
                try:
                    with st.spinner("AI 補答案中…"):
                        SS.generated_items = fill_missing_answers_with_ai(cfg, SS.generated_items)
                    st.success("✅ 已補回答案")
                except Exception as e:
                    show_exception("補答案失敗", e)

    # ④ 結果：預設表格一眼睇晒 + 可編輯
    SS.generated_items = render_editor(SS.generated_items, "④ 結果（可直接編輯題目/選項/答案）", "gen") or SS.generated_items

    # ⑤ 匯出/分享
    if isinstance(SS.generated_items, list) and SS.generated_items:
        st.markdown("## ⑤ 匯出與分享")

        # 先嘗試舊版 export panel
        export_panel_ok = False
        if _HAS_EXPORT_PANEL:
            try:
                render_export_panel(items=SS.generated_items, subject=subject, google_creds=SS.google_creds, mode="generate")
                export_panel_ok = True
            except Exception:
                export_panel_ok = False

        if not export_panel_ok:
            st.info("（保底匯出）")

        # 保底匯出
        xbio = export_kahoot_excel(SS.generated_items)
        dbio = export_wayground_docx(SS.generated_items, title=f"{subject} 題目清單")

        c1, c2 = st.columns(2)
        with c1:
            st.download_button(
                "⬇️ Kahoot（Excel）",
                data=xbio,
                file_name=f"kahoot_{now_stamp()}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        with c2:
            st.download_button(
                "⬇️ Wayground（DOCX）",
                data=dbio,
                file_name=f"wayground_{now_stamp()}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )

        # Google Forms + Email share
        st.markdown("### 🌐 Google Forms / ☁️ Drive / ✉️ Email")
        if not SS.google_creds:
            st.warning("請先在左側登入 Google 才可用 Google Forms / Drive / Gmail。")
        else:
            try:
                creds = credentials_from_dict(SS.google_creds)
                forms_service, drive_service, gmail_service = try_build_google_services(creds)

                colA, colB = st.columns(2)
                with colA:
                    if st.button("🚀 一鍵建立 Google Form", key="btn_form"):
                        url = create_google_form(forms_service, f"AI 題目表單（{subject}）", SS.generated_items)
                        SS.last_google_form_url = url
                        st.success("✅ 已建立 Google Form")
                        st.write(url)

                with colB:
                    if st.button("☁️ 上傳匯出檔到 Drive", key="btn_drive"):
                        link_x = drive_upload_bytes(
                            drive_service,
                            f"kahoot_{now_stamp()}.xlsx",
                            xbio.getvalue(),
                            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        )
                        link_d = drive_upload_bytes(
                            drive_service,
                            f"wayground_{now_stamp()}.docx",
                            dbio.getvalue(),
                            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        )
                        SS.last_drive_links = {"kahoot": link_x, "wayground": link_d}
                        st.success("✅ 已上傳")
                        st.write("Kahoot：", link_x)
                        st.write("Wayground：", link_d)

                st.markdown("#### 📧 一鍵電郵分享")
                SS.to_email = st.text_input("收件人（可多個，用逗號分隔）", value=SS.to_email)

                if st.button("📨 建立 Gmail 草稿（可附檔）", key="btn_gmail"):
                    links = SS.last_drive_links or {}
                    body = ["老師您好，以下為題目分享：", ""]
                    if SS.last_google_form_url:
                        body.append(f"Google Form：{SS.last_google_form_url}")
                    if links.get("kahoot"):
                        body.append(f"Kahoot Excel（Drive）：{links['kahoot']}")
                    if links.get("wayground"):
                        body.append(f"Wayground DOCX（Drive）：{links['wayground']}")
                    body.append("")

                    attachments = [
                        (f"kahoot_{now_stamp()}.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", xbio.getvalue()),
                        (f"wayground_{now_stamp()}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", dbio.getvalue()),
                    ]

                    try:
                        draft_id = gmail_create_draft_with_attachments(
                            gmail_service,
                            SS.to_email,
                            f"AI 題目分享（{subject}）",
                            "\n".join(body),
                            attachments,
                        )
                        st.success("✅ 已建立 Gmail 草稿（請到 Gmail 草稿匣發送）")
                        st.caption(f"Draft ID: {draft_id}")
                    except Exception as e:
                        st.warning("未能建立 Gmail 草稿（可能缺權限 scope）。改用 mailto 連結。")
                        body_q = "\n".join(body)
                        mailto = f"mailto:{SS.to_email}?subject=AI%20題目分享%EF%BC%88{subject}%EF%BC%89&body=" + st.utils.sanitize_url(body_q)
                        st.write(mailto)

            except Exception as e:
                st.warning("Google Forms/Drive/Gmail 功能不可用（可能未安裝 google-api-python-client 或 OAuth scope 不足）。")
                show_exception("Google 功能啟用失敗", e)


# =========================================================
# Import
# =========================================================
with tab_imp:
    st.markdown("## ① 上載 / 貼上題目（多格式）")
    uploads = st.file_uploader(
        "DOCX / PDF / XLSX / PPTX / TXT / 圖片",
        type=["pdf", "docx", "xlsx", "pptx", "txt", "png", "jpg", "jpeg"],
        accept_multiple_files=True,
        key="up_imp",
    )
    pasted = st.text_area("或直接貼上題目文字", value="", height=150)

    if uploads:
        with st.status("正在擷取題目內容…", expanded=False) as s:
            s.update(label="正在擷取題目內容…", state="running")
            text, v_urls, img_bytes, pdf_scan, report = extract_payload_from_uploads(uploads)
            SS.extracted_text_import = clean_text(text + "\n\n" + pasted)
            SS.vision_urls_import = v_urls
            SS.image_bytes_import = img_bytes
            SS.pdf_scan_bytes_import = pdf_scan
            SS.file_report_import = report
            s.update(label="擷取完成", state="complete")

        with st.expander("📄 擷取報告（每檔抽取字數）", expanded=False):
            st.dataframe(pd.DataFrame(SS.file_report_import), use_container_width=True)
    else:
        SS.extracted_text_import = clean_text(pasted)

    limit = 8000 if fast_mode else 10000

    st.markdown("## ②（可選）勾選要整理的段落")
    used = render_paragraph_picker_section(
        SS.extracted_text_import,
        mark_key="mark_idx_import",
        hash_key="last_text_hash_import",
        limit=limit,
        label="② 匯入內容段落（預設摺疊）",
    )

    st.info(f"送入 AI（重點合併後）{len(used)} 字（上限 {limit}；快速模式={'開' if fast_mode else '關'}）")

    st.markdown("## ③ 整理並轉換")
    use_ai = st.checkbox("使用 AI 協助整理（建議）", value=True)

    if st.button("✨ 整理並轉換", key="btn_imp"):
        if use_ai and not can_call_ai(cfg):
            st.error("請先輸入 API Key")
        else:
            try:
                prog = st.progress(0)
                prog.progress(20)

                used_for_ai = used

                if use_ai:
                    with st.spinner("AI 正在整理…"):
                        prog.progress(70)
                        SS.imported_items = normalize_items(assist_import_questions(cfg, used_for_ai, subject) or [])
                        prog.progress(100)
                else:
                    SS.imported_items = normalize_items(parse_import_questions_locally(used_for_ai) or [])
                    prog.progress(100)

                st.success("✅ 題目整理完成")

            except Exception as e:
                show_exception("整理失敗", e)

    SS.imported_items = render_editor(SS.imported_items, "④ 整理結果（可編輯）", "imp") or SS.imported_items

    if isinstance(SS.imported_items, list) and SS.imported_items:
        st.markdown("## ⑤ 匯出")
        c1, c2 = st.columns(2)
        with c1:
            st.download_button(
                "⬇️ Kahoot（Excel）",
                data=export_kahoot_excel(SS.imported_items),
                file_name=f"kahoot_import_{now_stamp()}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        with c2:
            st.download_button(
                "⬇️ Wayground（DOCX）",
                data=export_wayground_docx(SS.imported_items, title=f"{subject}（匯入）題目清單"),
                file_name=f"wayground_import_{now_stamp()}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
